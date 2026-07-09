from pathlib import Path
from PIL import Image
import io
import logging
import pytest
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from config import Config
from pipeline.extractor import extract, _is_svg, _downscale, _make_id


def test_extract_pptx_returns_records(sample_pptx):
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_pptx, config)
    assert len(records) == 1
    r = records[0]
    assert r.id  # non-empty hash
    assert "slide 1" in r.source_ref.lower()
    assert r.thumbnail_bytes
    assert r.recognition_bytes
    assert r.predicted_smiles is None
    assert r.is_chemical is None


def test_extract_pptx_thumbnail_size(sample_pptx):
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_pptx, config)
    img = Image.open(io.BytesIO(records[0].thumbnail_bytes))
    assert max(img.size) <= 64


def test_extract_pptx_recognition_size(sample_pptx):
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_pptx, config)
    img = Image.open(io.BytesIO(records[0].recognition_bytes))
    assert max(img.size) <= 128


def test_extract_docx_returns_records(sample_docx):
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_docx, config)
    assert len(records) == 1
    r = records[0]
    assert r.id
    assert r.thumbnail_bytes
    assert r.recognition_bytes


def test_extract_unsupported_format_raises(tmp_path):
    bad = tmp_path / "file.txt"
    bad.write_text("hello")
    with pytest.raises(ValueError, match="Unsupported"):
        extract(bad, Config())


# --- _is_svg ---

def test_is_svg_detects_svg_header():
    assert _is_svg(b"<svg xmlns='http://www.w3.org/2000/svg'><rect/></svg>")


def test_is_svg_detects_xml_then_svg():
    assert _is_svg(b"<?xml version='1.0'?><svg/>")


def test_is_svg_rejects_png():
    png_magic = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    assert not _is_svg(png_magic)


# --- _downscale transparency compositing ---

def test_downscale_rgba_composites_onto_white():
    # Semi-transparent blue on white → light blue, not black
    img = Image.new("RGBA", (100, 100), (0, 0, 255, 128))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    result = _downscale(buf.getvalue(), 64)
    out = Image.open(io.BytesIO(result))
    assert out.mode == "RGB"
    r, g, b = out.getpixel((50, 50))
    assert r > 100  # white background shows through: R should be well above 0


def test_downscale_palette_mode_produces_rgb():
    img = Image.new("P", (100, 100))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    result = _downscale(buf.getvalue(), 64)
    out = Image.open(io.BytesIO(result))
    assert out.mode == "RGB"


# --- _make_id ---

def test_make_id_is_deterministic_and_unique():
    assert _make_id(b"hello") == _make_id(b"hello")
    assert _make_id(b"hello") != _make_id(b"world")
    assert len(_make_id(b"hello")) == 64  # SHA-256 hex digest


# --- multi-image extraction ---

def test_extract_pptx_two_images(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(2):
        img = Image.new("RGB", (100, 100), color=(i * 100, 0, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(i + 1), Inches(1), Inches(1), Inches(1))
    path = tmp_path / "two.pptx"
    prs.save(str(path))
    records = extract(path, Config(thumbnail_max_size=64, recognition_max_size=64))
    assert len(records) == 2


def test_extract_docx_two_images_no_duplicates(tmp_path):
    doc = Document()
    for i in range(2):
        img = Image.new("RGB", (100, 100), color=(i * 100, 50, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        doc.add_picture(buf)
    path = tmp_path / "two.docx"
    doc.save(str(path))
    records = extract(path, Config(thumbnail_max_size=64, recognition_max_size=64))
    assert len(records) == 2


def test_extract_pptx_logs_opened_and_extracted(sample_pptx, caplog):
    caplog.set_level(logging.DEBUG, logger="pipeline.extractor")
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_pptx, config)
    messages = [r.message for r in caplog.records]
    assert any(m.startswith("Opened") and "1 images found" in m for m in messages)
    assert f"Extracted {records[0].source_ref}" in messages


def test_extract_docx_logs_opened_and_extracted(sample_docx, caplog):
    caplog.set_level(logging.DEBUG, logger="pipeline.extractor")
    config = Config(thumbnail_max_size=64, recognition_max_size=128)
    records = extract(sample_docx, config)
    messages = [r.message for r in caplog.records]
    assert any(m.startswith("Opened") and "1 images found" in m for m in messages)
    assert f"Extracted {records[0].source_ref}" in messages
