from pathlib import Path
from unittest.mock import patch
from pptx import Presentation
from docx import Document
from config import Config
from pipeline.extractor import extract
from pipeline.recognizer import recognize
from pipeline.reviewer import auto_accept
from pipeline.writer import write

_WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"


def test_full_pipeline_pptx(sample_pptx, tmp_path):
    config = Config(
        thumbnail_max_size=64,
        recognition_max_size=128,
        output_mode="new_file",
    )

    records = extract(sample_pptx, config)
    assert len(records) == 1

    with patch("pipeline.recognizer._run_decimer", return_value=("C1=CC=CC=C1", 0.95)):
        records = recognize(records, config)

    assert records[0].predicted_smiles == "C1=CC=CC=C1"
    assert records[0].recognition_bytes == b""

    records = auto_accept(records)
    assert records[0].approved_value == "C1=CC=CC=C1"
    assert records[0].is_chemical is True

    out_path = write(records, sample_pptx, config)
    assert out_path.exists()
    assert out_path.name == "sample_accessible.pptx"

    prs = Presentation(str(out_path))
    shape = list(prs.slides[0].shapes)[0]
    assert shape.element.nvPicPr.cNvPr.get("descr") == "C1=CC=CC=C1"


def test_full_pipeline_docx(sample_docx, tmp_path):
    config = Config(
        thumbnail_max_size=64,
        recognition_max_size=128,
        output_mode="new_file",
    )

    records = extract(sample_docx, config)
    assert len(records) == 1
    assert records[0].source_ref == "image 1"

    with patch("pipeline.recognizer._run_decimer", return_value=("C1=CC=CC=C1", 0.95)):
        records = recognize(records, config)

    assert records[0].predicted_smiles == "C1=CC=CC=C1"

    records = auto_accept(records)
    assert records[0].is_chemical is True

    out_path = write(records, sample_docx, config)
    assert out_path.exists()
    assert out_path.name == "sample_accessible.docx"

    doc = Document(str(out_path))
    doc_prs = doc.element.findall(f'.//{{{_WP}}}docPr')
    assert any(dp.get("descr") == "C1=CC=CC=C1" for dp in doc_prs)
