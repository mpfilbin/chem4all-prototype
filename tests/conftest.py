import io
from pathlib import Path
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from docx import Document


def _make_red_png(size=(200, 200)) -> bytes:
    img = Image.new("RGB", size, color=(255, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture()
def red_png_bytes() -> bytes:
    return _make_red_png()


@pytest.fixture()
def sample_pptx(tmp_path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_bytes = _make_red_png()
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    return path


@pytest.fixture()
def sample_docx(tmp_path) -> Path:
    doc = Document()
    img_bytes = _make_red_png()
    img_stream = io.BytesIO(img_bytes)
    doc.add_picture(img_stream)
    path = tmp_path / "sample.docx"
    doc.save(str(path))
    return path
