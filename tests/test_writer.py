import io
import logging
from pathlib import Path
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from config import Config
from models.image_record import ImageRecord
from pipeline.writer import write, _parse_source_ref_pptx


def _make_png() -> bytes:
    img = Image.new("RGB", (100, 100), color=(0, 255, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_pptx_with_image(tmp_path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_make_png()), Inches(1), Inches(1))
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


def _make_docx_with_image(tmp_path) -> Path:
    doc = Document()
    doc.add_picture(io.BytesIO(_make_png()))
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


def _approved_record(source_ref: str) -> ImageRecord:
    return ImageRecord(
        id="abc", source_ref=source_ref,
        thumbnail_bytes=b"", recognition_bytes=b"",
        predicted_smiles="C1=CC=CC=C1",
        approved_value="C1=CC=CC=C1",
        is_chemical=True,
    )


def test_writer_pptx_new_file(tmp_path):
    src = _make_pptx_with_image(tmp_path)
    record = _approved_record("slide 1, shape 1")
    out = write([record], src, Config(output_mode="new_file"))
    assert out != src
    assert out.exists()
    prs = Presentation(str(out))
    shape = list(prs.slides[0].shapes)[0]
    assert shape.name or True  # just verify file is readable


def test_writer_pptx_sets_alt_text(tmp_path):
    src = _make_pptx_with_image(tmp_path)
    record = _approved_record("slide 1, shape 1")
    out = write([record], src, Config(output_mode="new_file"))
    prs = Presentation(str(out))
    shape = list(prs.slides[0].shapes)[0]
    assert shape.element.nvPicPr.cNvPr.get("descr") == "C1=CC=CC=C1"


def test_writer_pptx_in_place(tmp_path):
    src = _make_pptx_with_image(tmp_path)
    record = _approved_record("slide 1, shape 1")
    out = write([record], src, Config(output_mode="in_place"))
    assert out == src


def test_parse_source_ref_pptx():
    assert _parse_source_ref_pptx("slide 1, shape 1") == (0, 0)
    assert _parse_source_ref_pptx("slide 3, shape 2") == (2, 1)


def test_writer_skips_non_chemical(tmp_path):
    src = _make_pptx_with_image(tmp_path)
    record = ImageRecord(
        id="abc", source_ref="slide 1, shape 1",
        thumbnail_bytes=b"", recognition_bytes=b"",
        is_chemical=False,
    )
    out = write([record], src, Config(output_mode="new_file"))
    prs = Presentation(str(out))
    shape = list(prs.slides[0].shapes)[0]
    # python-pptx sets descr="image.png" by default; our writer must not overwrite
    # it with a SMILES/chemical value for non-chemical records
    descr = shape.element.nvPicPr.cNvPr.get("descr")
    assert descr in (None, "image.png")  # unmodified default, no SMILES injected


# --- DOCX writer ---

_WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"


def _docx_doc_prs(path):
    return Document(str(path)).element.findall(f'.//{{{_WP}}}docPr')


def test_writer_docx_new_file(tmp_path):
    src = _make_docx_with_image(tmp_path)
    out = write([_approved_record("image 1")], src, Config(output_mode="new_file"))
    assert out != src
    assert out.exists()


def test_writer_docx_in_place(tmp_path):
    src = _make_docx_with_image(tmp_path)
    out = write([_approved_record("image 1")], src, Config(output_mode="in_place"))
    assert out == src


def test_writer_docx_sets_alt_text(tmp_path):
    src = _make_docx_with_image(tmp_path)
    out = write([_approved_record("image 1")], src, Config(output_mode="new_file"))
    assert any(dp.get("descr") == "C1=CC=CC=C1" for dp in _docx_doc_prs(out))


def test_writer_docx_skips_non_chemical(tmp_path):
    src = _make_docx_with_image(tmp_path)
    record = ImageRecord(
        id="abc", source_ref="image 1",
        thumbnail_bytes=b"", recognition_bytes=b"",
        is_chemical=False,
    )
    out = write([record], src, Config(output_mode="new_file"))
    assert not any(dp.get("descr") == "C1=CC=CC=C1" for dp in _docx_doc_prs(out))


def test_writer_pptx_logs_wrote(tmp_path, caplog):
    caplog.set_level(logging.DEBUG, logger="pipeline.writer")
    src = _make_pptx_with_image(tmp_path)
    record = _approved_record("slide 1, shape 1")
    out = write([record], src, Config(output_mode="new_file"))
    messages = [r.message for r in caplog.records]
    assert any(m.startswith(f"Wrote {out}") and "1 alt-texts applied" in m for m in messages)


def test_writer_pptx_no_wrote_log_when_nothing_approved(tmp_path, caplog):
    caplog.set_level(logging.DEBUG, logger="pipeline.writer")
    src = _make_pptx_with_image(tmp_path)
    record = ImageRecord(
        id="abc", source_ref="slide 1, shape 1",
        thumbnail_bytes=b"", recognition_bytes=b"",
        is_chemical=False,
    )
    write([record], src, Config(output_mode="new_file"))
    messages = [r.message for r in caplog.records]
    assert not any(m.startswith("Wrote") for m in messages)


def test_writer_docx_logs_wrote(tmp_path, caplog):
    caplog.set_level(logging.DEBUG, logger="pipeline.writer")
    src = _make_docx_with_image(tmp_path)
    out = write([_approved_record("image 1")], src, Config(output_mode="new_file"))
    messages = [r.message for r in caplog.records]
    assert any(m.startswith(f"Wrote {out}") and "1 alt-texts applied" in m for m in messages)
