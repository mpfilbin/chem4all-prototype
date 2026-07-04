from __future__ import annotations
import logging
import shutil
from pathlib import Path

from pptx import Presentation
from docx import Document

from config import Config
from models.image_record import ImageRecord

log = logging.getLogger(__name__)


def write(
    records: list[ImageRecord],
    source_path: Path,
    config: Config,
    output_path: Path | None = None,
) -> Path:
    if config.output_mode == "in_place" and output_path is None:
        dest = source_path
    elif output_path is not None:
        dest = output_path
    else:
        stem = source_path.stem
        dest = source_path.with_name(f"{stem}_accessible{source_path.suffix}")

    if dest != source_path:
        shutil.copy2(source_path, dest)

    suffix = source_path.suffix.lower()
    if suffix == ".pptx":
        _write_pptx(records, dest)
    elif suffix == ".docx":
        _write_docx(records, dest)
    return dest


def _parse_source_ref_pptx(source_ref: str) -> tuple[int, int]:
    # source_ref is 1-indexed; convert to 0-indexed
    parts = source_ref.replace(",", "").split()
    slide_idx = int(parts[1]) - 1
    shape_idx = int(parts[3]) - 1
    return slide_idx, shape_idx


def _write_pptx(records: list[ImageRecord], dest: Path) -> None:
    approved = [r for r in records if r.is_chemical is True and r.approved_value]
    if not approved:
        return
    prs = Presentation(str(dest))
    for record in approved:
        try:
            slide_idx, shape_idx = _parse_source_ref_pptx(record.source_ref)
            shape = list(prs.slides[slide_idx].shapes)[shape_idx]
            shape.element.nvPicPr.cNvPr.set("descr", record.approved_value)
        except Exception as exc:
            log.warning("Could not set alt-text for %s: %s", record.source_ref, exc)
    prs.save(str(dest))


def _write_docx(records: list[ImageRecord], dest: Path) -> None:
    approved = [r for r in records if r.is_chemical is True and r.approved_value]
    if not approved:
        return
    doc = Document(str(dest))

    # Rebuild the same image-index → rid mapping the extractor used
    seen: set[str] = set()
    image_idx = 0
    index_to_rid: dict[str, str] = {}
    for rid, rel in doc.part.rels.items():
        if "image" not in rel.reltype or rid in seen:
            continue
        seen.add(rid)
        image_idx += 1
        index_to_rid[f"image {image_idx}"] = rid

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"

    for record in approved:
        rid = index_to_rid.get(record.source_ref)
        if rid is None:
            log.warning("No relationship found for %s", record.source_ref)
            continue
        try:
            for blip in doc.element.findall(f'.//{{{A}}}blip[@{{{R}}}embed="{rid}"]'):
                node = blip.getparent()
                while node is not None:
                    if node.tag in (f'{{{WP}}}inline', f'{{{WP}}}anchor'):
                        break
                    node = node.getparent()
                if node is not None:
                    doc_pr = node.find(f'{{{WP}}}docPr')
                    if doc_pr is not None:
                        doc_pr.set("descr", record.approved_value)
        except Exception as exc:
            log.warning("Could not set alt-text for %s: %s", record.source_ref, exc)
    doc.save(str(dest))
