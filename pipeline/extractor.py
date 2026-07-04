from __future__ import annotations
import hashlib
import io
import logging
from collections.abc import Callable
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document

from config import Config
from models.image_record import ImageRecord

log = logging.getLogger(__name__)

_ProgressCallback = Callable[[int, int], None]


def extract(
    file_path: Path,
    config: Config,
    on_progress: _ProgressCallback | None = None,
) -> list[ImageRecord]:
    suffix = file_path.suffix.lower()
    if suffix == ".pptx":
        return _extract_pptx(file_path, config, on_progress)
    elif suffix == ".docx":
        return _extract_docx(file_path, config, on_progress)
    else:
        raise ValueError(f"Unsupported file format: {suffix}")


def _is_svg(raw: bytes) -> bool:
    head = raw[:512].lstrip()
    return head.startswith(b"<svg") or head.startswith(b"<?xml") or b"<svg" in head


def _svg_to_png(svg_bytes: bytes) -> bytes:
    try:
        import cairosvg
        return cairosvg.svg2png(bytestring=svg_bytes)
    except ImportError:
        raise OSError(
            "SVG image found but 'cairosvg' is not installed. "
            "Run: brew install cairo && pip install cairosvg"
        )


def _downscale(raw_bytes: bytes, max_size: int) -> bytes:
    if _is_svg(raw_bytes):
        raw_bytes = _svg_to_png(raw_bytes)
    img = Image.open(io.BytesIO(raw_bytes))
    if img.mode in ("RGBA", "LA", "P"):
        img = img.convert("RGBA")
        white = Image.new("RGBA", img.size, (255, 255, 255, 255))
        white.alpha_composite(img)
        img = white.convert("RGB")
    else:
        img = img.convert("RGB")
    img.thumbnail((max_size, max_size), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_id(raw_bytes: bytes) -> str:
    return hashlib.sha256(raw_bytes).hexdigest()


_PIC_TAG = "{http://schemas.openxmlformats.org/presentationml/2006/main}pic"


def _is_picture_shape(shape) -> bool:
    # MSO_SHAPE_TYPE.PICTURE covers standalone images; the p:pic tag also catches
    # images placed inside content placeholders, which python-pptx types as PLACEHOLDER.
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.element.tag == _PIC_TAG


def _extract_pptx(
    file_path: Path,
    config: Config,
    on_progress: _ProgressCallback | None = None,
) -> list[ImageRecord]:
    prs = Presentation(str(file_path))
    total = sum(
        1 for slide in prs.slides
        for shape in slide.shapes
        if _is_picture_shape(shape)
    )
    records: list[ImageRecord] = []
    extracted = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not _is_picture_shape(shape):
                continue
            try:
                raw = shape.image.blob
                records.append(ImageRecord(
                    id=_make_id(raw),
                    source_ref=f"slide {slide_idx}, shape {shape_idx}",
                    thumbnail_bytes=_downscale(raw, config.thumbnail_max_size),
                    recognition_bytes=_downscale(raw, config.recognition_max_size),
                ))
            except (OSError, AttributeError, ValueError) as exc:
                log.warning("Could not extract image slide %d shape %d: %s", slide_idx, shape_idx, exc)
            extracted += 1
            if on_progress:
                on_progress(extracted, total)
    return records


def _extract_docx(
    file_path: Path,
    config: Config,
    on_progress: _ProgressCallback | None = None,
) -> list[ImageRecord]:
    doc = Document(str(file_path))
    image_rids = [
        rid for rid, rel in doc.part.rels.items() if "image" in rel.reltype
    ]
    total = len(image_rids)
    records: list[ImageRecord] = []
    seen: set[str] = set()
    image_idx = 0
    for rid in image_rids:
        if rid in seen:
            continue
        seen.add(rid)
        rel = doc.part.rels[rid]
        try:
            raw = rel.target_part.blob
            image_idx += 1
            records.append(ImageRecord(
                id=_make_id(raw),
                source_ref=f"image {image_idx}",
                thumbnail_bytes=_downscale(raw, config.thumbnail_max_size),
                recognition_bytes=_downscale(raw, config.recognition_max_size),
            ))
        except (OSError, AttributeError, ValueError) as exc:
            log.warning("Could not extract image %d (rid=%s): %s", image_idx, rid, exc)
        if on_progress:
            on_progress(image_idx, total)
    return records
